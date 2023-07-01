import openai
import json
import time
import random
from pptx import Presentation
from io import BytesIO


def robot_print(text):
    for char in text:
        print(char, end="", flush=True)
        time.sleep(random.randrange(1, 5) / 100.0)
    print("\r")


def chatppt(topic: str, pages: int, api_key: str):
    language = "English"
    output_format = {
        "title": "example title",
        "pages": [
            {
                "title": "title for page 1",
                "subtitle": "subtitle for page 1",
                "content": [
                    {
                        "title": "title for bullet 1",
                        "desctription": "detail for bullet 1",
                    },
                    {
                        "title": "title for bullet 2",
                        "desctription": "detail for bullet 2",
                    },
                ],
            },
            {
                "title": "title for page 2",
                "subtitle": "subtitle for page 2",
                "content": [
                    {
                        "title": "title for bullet 1",
                        "desctription": "detail for bullet 1",
                    },
                    {
                        "title": "title for bullet 2",
                        "desctription": "detail for bullet 2",
                    },
                ],
            },
        ],
    }

    messages = [
        {
            "role": "user",
            "content": f"I'm going to prepare a presentation about {topic}, please help to outline detailed about this topic, output with JSON language with follow in format {output_format}, please help to generate {pages} pages, the bullet for each as much as possible and also write 2-3 lines for each bullet, please only return JSON format and use double quotes, please return the content in English",
        },
    ]
    robot_print(f"Your PPT will be generated in English")

    openai.api_key = api_key
    completion = openai.ChatCompletion.create(model="gpt-3.5-turbo", messages=messages)
    try:
        content = completion.choices[0].message.content
        # just replace ' to " is not a good soluation
        # print(content)
        content = json.loads(content.strip())
        return content
    except Exception as e:
        print("I'm a PPT assistant, your PPT generate failed, please retry later..")
        exit(1)
        # raise Exception("I'm PPT generate assistant, some error happened, please retry")


def generate_ppt(content: str, template):
    ppt = Presentation()
    if template:
        ppt = Presentation(template)
    first_slide_layout = ppt.slide_layouts[0]

    # """ Ref for slide types:
    # 0 ->  title and subtitle
    # 1 ->  title and content
    # 2 ->  section header
    # 3 ->  two content
    # 4 ->  Comparison
    # 5 ->  Title only
    # 6 ->  Blank
    # 7 ->  Content with caption
    # 8 ->  Pic with caption
    # """

    slide = ppt.slides.add_slide(first_slide_layout)
    slide.shapes.title.text = content.get("title", "")
    slide.placeholders[1].text = "Generate by GPTPPT"
    pages = content.get("pages", [])
    robot_print(f"Your PPT have {len(pages)} pages.")
    for i, page in enumerate(pages):
        page_title = page.get("title", "")
        robot_print(f"page {i + 1}: {page_title}")
        bullet_layout = ppt.slide_layouts[1]
        bullet_slide = ppt.slides.add_slide(bullet_layout)
        bullet_spahe = bullet_slide.shapes
        bullet_slide.shapes.title.text = page_title

        body_shape = bullet_spahe.placeholders[1]
        for bullet in page.get("content", []):
            paragraph = body_shape.text_frame.add_paragraph()
            paragraph.text = bullet.get("title", "")
            paragraph.level = 1

            paragraph = body_shape.text_frame.add_paragraph()
            paragraph.text = bullet.get("description", "")
            paragraph.level = 2

    ppt_name = content.get("title", "")
    ppt_name = f"{ppt_name}.pptx"
    binary_output = BytesIO()
    ppt.save(binary_output)
    robot_print("Generate done, enjoy!")
    return binary_output


def main(topic: str, pages: int, api_key: str, template_path):
    ppt_content = chatppt(topic, pages, api_key)
    binary_output = generate_ppt(ppt_content, template=template_path)
    return binary_output
